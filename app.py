import os, json, io, re, requests, httpx
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import anthropic

app = Flask(__name__)
CORS(app)

# ── Keys from environment (set these in Railway) ─────────────
CLAUDE_KEY  = os.environ.get("CLAUDE_API_KEY", "")
SERPER_KEY  = os.environ.get("SERPER_API_KEY", "")
YOUTUBE_KEY = os.environ.get("YOUTUBE_API_KEY", "")

# ── Colors ───────────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0A, 0x14)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0x7C, 0x6A, 0xF7)
ACCENT2  = RGBColor(0x4F, 0xD1, 0xC5)
MUTED    = RGBColor(0x88, 0x88, 0xA8)
CARD     = RGBColor(0x1A, 0x1A, 0x2E)
GREEN    = RGBColor(0x1E, 0x84, 0x49)
RED      = RGBColor(0xC0, 0x39, 0x2B)
AMBER    = RGBColor(0xD4, 0xA0, 0x17)
TEXT     = RGBColor(0xE8, 0xE8, 0xF0)

SLD_W = Inches(13.33)
SLD_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────
def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def slide_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or BG

def add_eyebrow(slide, text):
    add_text(slide, text, 0.4, 0.2, 12, 0.3, size=8, bold=True, color=ACCENT2)

def add_title(slide, text):
    add_text(slide, text, 0.4, 0.5, 12, 0.6, size=26, bold=True, color=WHITE)

def add_divider(slide, y=1.2):
    add_rect(slide, 0.4, y, 12.4, 0.02, rgb(0x2A,0x2A,0x3A))

def tag_color(tag):
    t = tag.upper()
    if "POSITIVE" in t and "MIXED" not in t: return GREEN, WHITE
    if "STRENGTH" in t: return GREEN, WHITE
    if "NEGATIVE" in t or "CRITICAL" in t or "HIGH" in t: return RED, WHITE
    if "MIXED" in t or "MANAGE" in t or "MEDIUM" in t or "AMBER" in t: return AMBER, rgb(0x1A,0x1A,0x2E)
    if "LEVERAGE" in t or "INFO" in t: return ACCENT, WHITE
    return rgb(0x55,0x55,0x70), WHITE

# ── Serper search (runs 35+ queries) ─────────────────────────
def search(query):
    try:
        r = requests.post("https://google.serper.dev/search",
            headers={"X-API-KEY": SERPER_KEY, "Content-Type": "application/json"},
            json={"q": query, "num": 5}, timeout=8)
        data = r.json()
        results = [f"• {x['title']}: {x['snippet']}" for x in data.get("organic", [])]
        return "\n".join(results)[:400] or "No results"
    except: return "Search unavailable"

def scrape_website_socials(website):
    """Visit brand website and extract real social media links"""
    if not website:
        return {}
    try:
        import re
        url = website if website.startswith("http") else f"https://{website}"
        headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
        r = requests.get(url, headers=headers, timeout=12, allow_redirects=True)
        html = r.text
        found = {}
        patterns = {
            "instagram": r"(?:href=["\'])?(https?://)?(?:www\.)?instagram\.com/([A-Za-z0-9_.]+)",
            "facebook":  r"(?:href=["\'])?(https?://)?(?:www\.)?facebook\.com/([A-Za-z0-9_./-]+)",
            "youtube":   r"(?:href=["\'])?(https?://)?(?:www\.)?youtube\.com/(?:@|c/|channel/)?([A-Za-z0-9_-]+)",
            "linkedin":  r"(?:href=["\'])?(https?://)?(?:www\.)?linkedin\.com/(?:company/|in/)?([A-Za-z0-9_-]+)",
            "twitter":   r"(?:href=["\'])?(https?://)?(?:www\.)?(?:twitter|x)\.com/([A-Za-z0-9_]+)",
        }
        for platform, pattern in patterns.items():
            match = re.search(pattern, html, re.IGNORECASE)
            if match:
                handle = match.group(2) if match.lastindex >= 2 else ""
                if handle and handle not in ["share","sharer","intent","home","login","signup"]:
                    found[platform] = handle
        print(f"[WEBSITE] Found social links: {found}")
        return found
    except Exception as e:
        print(f"[WEBSITE] Scrape failed: {e}")
        return {}

def deep_research(brand, website=""):
    import concurrent.futures

    # Clean domain for targeted searches
    domain = website.replace("https://","").replace("http://","").replace("www.","").strip("/") if website else ""

    # Get real social handles from website
    social = scrape_website_socials(website) if website else {}
    ig  = social.get("instagram", "")
    fb  = social.get("facebook", "")
    yt  = social.get("youtube", "")
    li  = social.get("linkedin", "")
    tw  = social.get("twitter", "")

    # Build targeted queries using real handles where found
    ig_q  = f"instagram.com/{ig} followers engagement posts reels" if ig else f"{brand} instagram followers engagement posts 2025"
    fb_q  = f"facebook.com/{fb} page followers likes posts" if fb else f"{brand} facebook page followers likes engagement"
    yt_q  = f"youtube.com/{yt} subscribers views videos" if yt else f"{brand} youtube channel subscribers views"
    li_q  = f"linkedin.com/company/{li} employees followers" if li else f"{brand} linkedin company followers employees"
    tw_q  = f"twitter.com/{tw} followers engagement" if tw else f"{brand} twitter X followers engagement"

    # SEMrush deep queries - using real domain if available
    sem1  = f"semrush.com {domain} monthly visits organic traffic keywords" if domain else f"semrush {brand} traffic keywords overview"
    sem2  = f"{domain} top organic keywords search volume ranking" if domain else f"{brand} top keywords search volume SEO"
    sem3  = f"{domain} backlinks referring domains authority score" if domain else f"{brand} backlinks domain authority"

    # Meta Ad Library deep queries
    meta1 = f'facebook.com/ads/library {brand} active ads India'
    meta2 = f'{brand} facebook instagram ads running creative format video image'
    meta3 = f'{brand} meta ads CTA hook messaging spend range date started'

    # Google Ad Transparency deep queries
    gads1 = f'adstransparency.google.com {brand} ads India'
    gads2 = f'{brand} google search ads display ads youtube ads running'
    gads3 = f'{brand} google ads keywords bidding strategy ppc campaigns'

    queries = [
        # Company & financial (7)
        f"{brand} brand overview company history founding India",
        f"{brand} revenue funding investors valuation 2025",
        f"{brand} company registration MCA CIN directors India",
        f"{brand} startup India DPIIT recognition certificate",
        f"{brand} zauba tofler company financial data",
        f"{brand} inc42 yourStory funding rounds investors",
        f"{brand} annual report profit loss revenue growth",
        # Social media (5)
        ig_q, fb_q, yt_q, li_q, tw_q,
        # SEMrush deep (3)
        sem1, sem2, sem3,
        # Meta Ads deep (3)
        meta1, meta2, meta3,
        # Google Ads deep (3)
        gads1, gads2, gads3,
        # Market & competition (4)
        f"{brand} target audience customer profile demographics India",
        f"{brand} competitors analysis market share India",
        f"{brand} market size TAM total addressable market India",
        f"{brand} SWOT analysis strengths weaknesses opportunities",
        # Reviews (5)
        f"{brand} site:amazon.in reviews rating",
        f"{brand} site:flipkart.com reviews rating",
        f"{brand} site:trustpilot.com reviews",
        f"{brand} google reviews rating customer feedback",
        f"{brand} site:reddit.com OR site:quora.com reviews discussion",
        # Customer voice (3)
        f"{brand} customers love positive reviews testimonials",
        f"{brand} complaints problems negative reviews issues",
        f"{brand} pricing value for money affordable expensive",
        # Press & marketing (4)
        f"{brand} press coverage Economic Times Forbes India media",
        f"{brand} marketing strategy campaigns advertising 2025",
        f"{brand} brand positioning tagline unique selling point",
        f"{brand} influencer celebrity ambassador collaboration",
    ]

    keys = [
        "overview","financials","registration","startupIndia","zauba","funding","revenue",
        "instagram","facebook","youtube","linkedin","twitter",
        "semrush1","semrush2","semrush3",
        "metaAds1","metaAds2","metaAds3",
        "googleAds1","googleAds2","googleAds3",
        "audience","competitors","tam","swot",
        "amazon","flipkart","trustpilot","googleReviews","reddit",
        "positives","negatives","pricing",
        "press","marketing","positioning","influencers",
    ]

    with concurrent.futures.ThreadPoolExecutor(max_workers=12) as ex:
        results = list(ex.map(search, queries))
    raw = dict(zip(keys, results))
    raw["social_handles"] = social
    raw["domain"] = domain
    return raw

# ── YouTube real data ─────────────────────────────────────────
def get_youtube(brand):
    if not YOUTUBE_KEY: return None
    try:
        s = requests.get("https://www.googleapis.com/youtube/v3/search",
            params={"part":"snippet","q":f"{brand} official","type":"channel",
                    "maxResults":1,"key":YOUTUBE_KEY}, timeout=8).json()
        if not s.get("items"): return None
        cid = s["items"][0]["id"]["channelId"]
        c = requests.get("https://www.googleapis.com/youtube/v3/channels",
            params={"part":"statistics,snippet","id":cid,"key":YOUTUBE_KEY}, timeout=8).json()
        if not c.get("items"): return None
        ch = c["items"][0]
        st = ch["statistics"]
        return {
            "name": ch["snippet"]["title"],
            "subscribers": f"{int(st.get('subscriberCount',0)):,}",
            "views": f"{int(st.get('viewCount',0)):,}",
            "videos": st.get("videoCount","0")
        }
    except: return None

# ── Claude analysis ───────────────────────────────────────────
def claude_analysis(brand, raw, yt):
    client = anthropic.Anthropic(api_key=CLAUDE_KEY, http_client=httpx.Client())
    yt_info = f"REAL YouTube: {yt['subscribers']} subscribers, {yt['views']} views, {yt['videos']} videos" if yt else "YouTube: estimate from research"

    prompt = f"""Brand intelligence report for "{brand}". Return ONLY valid JSON.
NEVER say "Data not available" - always estimate with (est.) label.
Keep ALL string values under 100 characters.

{yt_info}

DATA:
overview: {raw['overview'][:200]}
financials: {raw['financials'][:150]}
registration: {raw['registration'][:120]}
instagram: {raw['instagram'][:100]}
facebook: {raw['facebook'][:100]}
youtube: {raw['youtube'][:100]}
metaAds: {raw['metaAds'][:120]}
semrush: {raw['semrush'][:100]}
audience: {raw['audience'][:150]}
competitors: {raw['competitors'][:150]}
amazon: {raw['amazon'][:120]}
flipkart: {raw['flipkart'][:100]}
reddit: {raw['reddit'][:100]}
quora: {raw['quora'][:100]}
positives: {raw['positives'][:120]}
negatives: {raw['negatives'][:120]}
press: {raw['press'][:120]}
funding: {raw['funding'][:120]}
tam: {raw['tam'][:100]}
swot: {raw['swot'][:120]}
pricing: {raw['pricing'][:100]}
quickcommerce: {raw['quickcommerce'][:100]}
brandStory: {raw['brandStory'][:150]}

Fill this JSON completely. Every field must have a real or estimated value:
{{
  "brandName": "",
  "tagline": "",
  "category": "",
  "platform_footprint": {{
    "platforms": [
      {{"name":"Instagram","handle":"","followers":"","posts":"","status":""}},
      {{"name":"Facebook","handle":"","followers":"","posts":"","status":""}},
      {{"name":"YouTube","handle":"","followers":"","posts":"","status":""}},
      {{"name":"LinkedIn","handle":"","followers":"","posts":"","status":""}},
      {{"name":"X (Twitter)","handle":"","followers":"","posts":"","status":""}}
    ],
    "insight": ""
  }},
  "youtube_real": {{"subscribers":"","total_views":"","video_count":"","top_content":"","gap":""}},
  "engagement": {{"ig_rate":"","ig_likes":"","ig_comments":"","insight":""}},
  "sentiment": {{
    "positive_pct": 65,
    "mixed_pct": 20,
    "negative_pct": 15,
    "themes": [
      {{"theme":"","signal":"","tag":"STRENGTH"}},
      {{"theme":"","signal":"","tag":"CRITICAL"}},
      {{"theme":"","signal":"","tag":"MANAGE"}},
      {{"theme":"","signal":"","tag":"LEVERAGE"}},
      {{"theme":"","signal":"","tag":"WHITESPACE"}}
    ]
  }},
  "financial": {{
    "total_funding":"","last_round":"","investors":["","",""],
    "revenue":"","valuation":"","employees":"","founded":"","hq":""
  }},
  "company_legal": {{
    "cin":"","registered_name":"","incorporation_date":"",
    "directors":["",""],"dpiit_status":"","startup_india":"","mca_status":""
  }},
  "meta_ads": {{
    "total_ads":"","formats":[
      {{"type":"Video Ads","pct":70}},
      {{"type":"Static Image","pct":20}},
      {{"type":"Carousel","pct":10}}
    ],
    "top_hook":"","gap":""
  }},
  "semrush_data": {{
    "monthly_visits":"",
    "organic_keywords":"",
    "domain_rating":"",
    "backlinks":"",
    "traffic_trend":"Growing/Flat/Declining",
    "top_keywords":["","","","",""]
  }},
  "meta_ads_data": {{
    "total_ads":"",
    "primary_format":"",
    "primary_cta":"",
    "avg_duration":"",
    "formats":[
      {{"type":"Video Ads","pct":70}},
      {{"type":"Static Image","pct":20}},
      {{"type":"Carousel","pct":10}}
    ],
    "top_hooks":["",""],
    "gap":""
  }},
  "google_ads_data": {{
    "total_ads":"",
    "ad_types":"",
    "primary_keywords":"",
    "est_spend":"",
    "ad_type_breakdown":[
      {{"type":"Search Ads","pct":50,"note":""}},
      {{"type":"Display Ads","pct":30,"note":""}},
      {{"type":"YouTube Ads","pct":20,"note":""}}
    ],
    "insight":""
  }},
  "icp": {{
    "icps": [
      {{"name":"","score":"9/10","pct":"~65%","who":"","pain":"","awareness":"","fit":"STRONG FIT"}},
      {{"name":"","score":"7/10","pct":"~25%","who":"","pain":"","awareness":"","fit":"PARTIAL"}},
      {{"name":"","score":"5/10","pct":"~10%","who":"","pain":"","awareness":"","fit":"EMERGING"}}
    ]
  }},
  "hooks": [
    {{"id":"H1","name":"","score":"8/10","pct":"35%","structure":"","insight":""}},
    {{"id":"H2","name":"","score":"7/10","pct":"25%","structure":"","insight":""}},
    {{"id":"H3","name":"","score":"6/10","pct":"20%","structure":"","insight":""}},
    {{"id":"H4","name":"","score":"5/10","pct":"15%","structure":"","insight":""}}
  ],
  "conversation": [
    {{"platform":"Amazon Reviews","summary":"","tag":"MIXED-POSITIVE"}},
    {{"platform":"Flipkart","summary":"","tag":"POSITIVE"}},
    {{"platform":"Trustpilot","summary":"","tag":"MIXED"}},
    {{"platform":"Google Reviews","summary":"","tag":"POSITIVE"}},
    {{"platform":"Reddit","summary":"","tag":"MIXED-NEGATIVE"}},
    {{"platform":"Quora","summary":"","tag":"MIXED"}},
    {{"platform":"Quick Commerce","summary":"","tag":"MIXED"}},
    {{"platform":"Business Press","summary":"","tag":"POSITIVE"}}
  ],
  "ratings": [
    {{"platform":"Brand Website","rating":"4.X","reviews":"","tag":"POSITIVE"}},
    {{"platform":"Amazon.in","rating":"4.X","reviews":"","tag":"MIXED-POSITIVE"}},
    {{"platform":"Flipkart","rating":"4.X","reviews":"","tag":"POSITIVE"}},
    {{"platform":"Google","rating":"4.X","reviews":"","tag":"POSITIVE"}},
    {{"platform":"Trustpilot","rating":"3.X","reviews":"","tag":"LOW SIGNAL"}}
  ],
  "competitors": {{
    "list": [
      {{"name":"","position":"","strength":"","gap":""}},
      {{"name":"","position":"","strength":"","gap":""}},
      {{"name":"","position":"","strength":"","gap":""}}
    ],
    "key_gap": ""
  }},
  "voice_of_customer": {{
    "loves": [
      {{"theme":"","quote":""}},{{"theme":"","quote":""}},
      {{"theme":"","quote":""}},{{"theme":"","quote":""}}
    ],
    "frustrations": [
      {{"theme":"","quote":""}},{{"theme":"","quote":""}},
      {{"theme":"","quote":""}},{{"theme":"","quote":""}}
    ]
  }},
  "press": {{
    "drivers": [
      {{"title":"","detail":""}},{{"title":"","detail":""}},
      {{"title":"","detail":""}},{{"title":"","detail":""}}
    ],
    "sentiment": ""
  }},
  "gaps": [
    {{"title":"","priority":"HIGH","description":""}},
    {{"title":"","priority":"HIGH","description":""}},
    {{"title":"","priority":"MEDIUM","description":""}},
    {{"title":"","priority":"MEDIUM","description":""}},
    {{"title":"","priority":"MEDIUM","description":""}}
  ],
  "recommendations": [
    {{"number":"01","title":"","description":""}},
    {{"number":"02","title":"","description":""}},
    {{"number":"03","title":"","description":""}},
    {{"number":"04","title":"","description":""}},
    {{"number":"05","title":"","description":""}}
  ]
}}"""

    # Use streaming to avoid timeout
    full_text = ""
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=6000,
        messages=[{"role":"user","content":prompt}]
    ) as stream:
        for text in stream.text_stream:
            full_text += text

    clean = re.sub(r'```json|```','', full_text).strip()
    try:
        return json.loads(clean)
    except:
        # Repair truncated JSON
        ob = clean.count('{') - clean.count('}')
        ob2 = clean.count('[') - clean.count(']')
        clean = re.sub(r',?\s*"[^"]*$','', clean).rstrip(',').rstrip()
        for _ in range(ob2): clean += ']'
        for _ in range(ob): clean += '}'
        return json.loads(clean)

# ── PPT Builder ───────────────────────────────────────────────
def build_ppt(d, yt_real):
    prs = Presentation()
    prs.slide_width  = SLD_W
    prs.slide_height = SLD_H
    blank = prs.slide_layouts[6]

    def new_slide():
        s = prs.slides.add_slide(blank)
        slide_bg(s)
        add_rect(s, 0, 0, 0.07, 7.5, ACCENT)
        return s

    # ── COVER ────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    slide_bg(cover, rgb(0x0A,0x0A,0x14))
    add_rect(cover, 0, 0, 0.12, 7.5, ACCENT)
    add_rect(cover, 0, 0, 13.33, 7.5, ACCENT)
    # Re-draw dark bg on top of accent
    add_rect(cover, 0.12, 0, 13.21, 7.5, BG)
    name = d.get("brandName", "Brand").upper()
    add_text(cover, name, 0.5, 1.2, 12, 1.5, size=54, bold=True, color=WHITE)
    add_text(cover, d.get("tagline","Brand Intelligence Report"), 0.5, 2.8, 11, 0.6, size=18, italic=True, color=ACCENT2)
    add_rect(cover, 0.5, 3.3, 3, 0.04, ACCENT2)
    add_text(cover, "BRAND INTELLIGENCE REPORT", 0.5, 3.5, 4, 0.4, size=9, bold=True, color=ACCENT)
    add_text(cover, d.get("category","") + "  ·  June 2026  ·  Automated Research", 0.5, 6.8, 10, 0.4, size=9, color=MUTED)

    # ── SLIDE 1: Platform Footprint ──────────────────────────
    s1 = new_slide()
    add_eyebrow(s1, "WHERE THE BRAND SHOWS UP")
    add_title(s1, "Platform Footprint at a Glance")
    add_divider(s1)
    platforms = d.get("platform_footprint",{}).get("platforms",[])
    col_w = 12.0 / max(len(platforms),1)
    for i, p in enumerate(platforms):
        cx = 0.5 + i * col_w
        add_rect(s1, cx, 1.4, col_w-0.15, 2.9, CARD)
        st = p.get("status","")
        sc = {"Primary channel": GREEN, "Secondary": ACCENT,
              "B2B": rgb(0x0F,0x6E,0x56), "Under-leveraged": AMBER, "Inactive": RED}.get(st, MUTED)
        add_rect(s1, cx, 1.4+2.9-0.38, col_w-0.15, 0.38, sc)
        add_text(s1, st, cx, 1.4+2.52, col_w-0.15, 0.35, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s1, p.get("name",""), cx+0.1, 1.55, col_w-0.3, 0.35, size=12, bold=True, color=WHITE)
        add_text(s1, p.get("handle",""), cx+0.1, 1.88, col_w-0.3, 0.28, size=9, color=MUTED)
        add_text(s1, p.get("followers","—"), cx+0.1, 2.2, col_w-0.3, 0.55, size=22, bold=True, color=ACCENT)
        add_text(s1, "followers", cx+0.1, 2.72, col_w-0.3, 0.28, size=9, color=MUTED)
        add_text(s1, p.get("posts",""), cx+0.1, 3.0, col_w-0.3, 0.28, size=9, color=TEXT)
    insight = d.get("platform_footprint",{}).get("insight","")
    if insight:
        add_rect(s1, 0.5, 4.5, 12.3, 0.65, rgb(0x1A,0x1A,0x2E))
        add_rect(s1, 0.5, 4.5, 0.07, 0.65, ACCENT)
        add_text(s1, insight, 0.7, 4.57, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE 2: Engagement ──────────────────────────────────
    s2 = new_slide()
    add_eyebrow(s2, "HOW THE AUDIENCE BEHAVES")
    add_title(s2, "Audience & Engagement")
    add_divider(s2)
    eng = d.get("engagement",{})
    metrics = [
        ("Avg. Engagement Rate (IG)", eng.get("ig_rate","N/A"), "vs industry benchmark"),
        ("Avg. Likes per Post",       eng.get("ig_likes","N/A"), "Reaction volume signal"),
        ("Avg. Comments per Post",    eng.get("ig_comments","N/A"), "Conversation depth"),
    ]
    for i, (label, val, note) in enumerate(metrics):
        my = 1.45 + i * 1.1
        add_rect(s2, 0.5, my, 12.3, 0.95, CARD)
        add_text(s2, val, 1.4, my+0.1, 2.5, 0.6, size=28, bold=True, color=ACCENT)
        add_text(s2, label, 3.5, my+0.08, 5, 0.38, size=13, bold=True, color=WHITE)
        add_text(s2, note, 3.5, my+0.5, 8.5, 0.3, size=10, color=MUTED)
    insight2 = eng.get("insight","")
    if insight2:
        add_rect(s2, 0.5, 4.9, 12.3, 0.65, rgb(0x0A,0x1A,0x30))
        add_rect(s2, 0.5, 4.9, 0.07, 0.65, ACCENT)
        add_text(s2, insight2, 0.7, 4.97, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE 3: Sentiment ───────────────────────────────────
    s3 = new_slide()
    add_eyebrow(s3, "WHAT THE TONE LOOKS LIKE")
    add_title(s3, "Sentiment Overview")
    add_divider(s3)
    sent = d.get("sentiment",{})
    pos = sent.get("positive_pct",65)
    mix = sent.get("mixed_pct",20)
    neg = sent.get("negative_pct",15)
    bar_w = 5.5
    add_rect(s3, 0.5, 1.5, bar_w*(pos/100), 0.55, GREEN)
    add_rect(s3, 0.5+bar_w*(pos/100), 1.5, bar_w*(mix/100), 0.55, AMBER)
    add_rect(s3, 0.5+bar_w*((pos+mix)/100), 1.5, bar_w*(neg/100), 0.55, RED)
    add_text(s3, f"{pos}%", 0.5, 2.12, 1.5, 0.35, size=14, bold=True, color=GREEN)
    add_text(s3, "Positive", 0.5, 2.45, 1.5, 0.28, size=9, color=MUTED)
    add_text(s3, f"{mix}%", 0.5+bar_w*(pos/100), 2.12, 1.5, 0.35, size=14, bold=True, color=AMBER)
    add_text(s3, "Mixed", 0.5+bar_w*(pos/100), 2.45, 1.5, 0.28, size=9, color=MUTED)
    add_text(s3, f"{neg}%", 0.5+bar_w*((pos+mix)/100), 2.12, 1.5, 0.35, size=14, bold=True, color=RED)
    add_text(s3, "Negative", 0.5+bar_w*((pos+mix)/100), 2.45, 1.5, 0.28, size=9, color=MUTED)
    add_text(s3, f"{pos}%", 7.0, 1.3, 2.5, 1.1, size=60, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s3, "net positive", 7.0, 2.35, 2.5, 0.3, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    themes = sent.get("themes",[])
    for i, t in enumerate(themes):
        ty = 2.9 + i * 0.45
        add_rect(s3, 0.5, ty, 12.3, 0.38, CARD)
        add_text(s3, t.get("theme",""), 0.65, ty+0.07, 3, 0.26, size=11, bold=True, color=WHITE)
        add_text(s3, t.get("signal",""), 3.9, ty+0.07, 7.2, 0.26, size=10, color=TEXT)
        tc, tfc = tag_color(t.get("tag",""))
        add_rect(s3, 11.35, ty+0.07, 1.3, 0.24, tc)
        add_text(s3, t.get("tag",""), 11.35, ty+0.07, 1.3, 0.24, size=7, bold=True, color=tfc, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: Meta Ads ────────────────────────────────────
    s4 = new_slide()
    add_eyebrow(s4, "WHAT THEY RUN ON META")
    add_title(s4, "Ad Library Overview")
    add_divider(s4)
    ads = d.get("meta_ads",{})
    stat_labels = [("Total Ads", ads.get("total_ads","N/A")), ("Top Hook", ads.get("top_hook","N/A")), ("Platforms", "Facebook, Instagram")]
    for i, (lbl, val) in enumerate(stat_labels):
        mx = 0.5 + i * 4.1
        add_rect(s4, mx, 1.38, 3.9, 0.95, CARD)
        add_text(s4, lbl, mx+0.15, 1.45, 3.6, 0.28, size=8, bold=True, color=MUTED)
        add_text(s4, val, mx+0.15, 1.72, 3.6, 0.52, size=11, bold=True, color=WHITE, wrap=True)
    add_text(s4, "AD FORMAT BREAKDOWN", 0.5, 2.55, 6, 0.3, size=8, bold=True, color=MUTED)
    formats = ads.get("formats",[])
    for i, f in enumerate(formats):
        fy = 2.92 + i * 0.72
        add_text(s4, f.get("type",""), 0.5, fy, 2.5, 0.28, size=11, bold=True, color=WHITE)
        pct = f.get("pct",0)
        add_rect(s4, 3.1, fy+0.04, 5.5, 0.32, rgb(0x2A,0x2A,0x3A))
        add_rect(s4, 3.1, fy+0.04, 5.5*(pct/100), 0.32, ACCENT)
        add_text(s4, f"{pct}%", 8.8, fy+0.04, 0.8, 0.32, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    gap = ads.get("gap","")
    if gap:
        add_rect(s4, 0.5, 5.1, 12.3, 0.65, rgb(0x2A,0x0A,0x0A))
        add_rect(s4, 0.5, 5.1, 0.07, 0.65, RED)
        add_text(s4, f"CRITICAL GAP — {gap}", 0.7, 5.17, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE 5: ICP ─────────────────────────────────────────
    s5 = new_slide()
    add_eyebrow(s5, "WHO THEY TARGET — AND THE GAPS")
    add_title(s5, "ICP Deep Dive")
    add_divider(s5)
    icps = d.get("icp",{}).get("icps",[])
    col = 12.3 / max(len(icps),1)
    fit_colors = {"STRONG FIT": GREEN, "STRONGEST FIT": GREEN, "PARTIAL": AMBER, "EMERGING": ACCENT, "WEAK FIT": RED, "IGNORED": RED}
    for i, icp in enumerate(icps):
        ix = 0.5 + i * col
        add_rect(s5, ix, 1.38, col-0.2, 3.8, CARD)
        add_rect(s5, ix+0.15, 1.52, 0.7, 0.35, ACCENT)
        add_text(s5, icp.get("score",""), ix+0.15, 1.52, 0.7, 0.35, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s5, icp.get("name",""), ix+0.95, 1.55, col-1.2, 0.32, size=12, bold=True, color=WHITE)
        add_text(s5, icp.get("pct",""), ix+0.15, 1.88, col-0.35, 0.25, size=9, italic=True, color=MUTED)
        add_rect(s5, ix+0.15, 2.18, col-0.35, 0.02, rgb(0x2A,0x2A,0x3A))
        for j, (lbl, key) in enumerate([("WHO THEY ARE","who"),("CORE PAIN","pain"),("AWARENESS","awareness")]):
            sy = 2.25 + j * 0.72
            add_text(s5, lbl, ix+0.15, sy, col-0.35, 0.25, size=7, bold=True, color=MUTED)
            add_text(s5, icp.get(key,""), ix+0.15, sy+0.24, col-0.35, 0.42, size=10, color=TEXT, wrap=True)
        fc = fit_colors.get(icp.get("fit","PARTIAL"), AMBER)
        add_rect(s5, ix+0.15, 4.72, col-0.35, 0.32, fc)
        add_text(s5, icp.get("fit",""), ix+0.15, 4.72, col-0.35, 0.32, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── SLIDE 6: Hooks ───────────────────────────────────────
    s6 = new_slide()
    add_eyebrow(s6, "HOW THEY COMMUNICATE")
    add_title(s6, "Hook & Messaging Patterns")
    add_divider(s6)
    hooks = d.get("hooks",[])
    for i, h in enumerate(hooks):
        hy = 1.38 + i * 1.0
        add_rect(s6, 0.5, hy, 12.3, 0.88, CARD)
        add_rect(s6, 0.6, hy+0.18, 0.55, 0.52, AMBER)
        add_text(s6, h.get("id",""), 0.6, hy+0.18, 0.55, 0.52, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s6, h.get("name",""), 1.28, hy+0.1, 3, 0.35, size=13, bold=True, color=WHITE)
        add_text(s6, f"{h.get('score','')}  ·  {h.get('pct','')}", 1.28, hy+0.48, 3, 0.28, size=9, color=MUTED)
        add_text(s6, "STRUCTURE:", 4.5, hy+0.1, 1.1, 0.25, size=7, bold=True, color=MUTED)
        add_text(s6, h.get("structure",""), 5.65, hy+0.08, 3.8, 0.38, size=10, italic=True, color=TEXT, wrap=True)
        add_text(s6, "INSIGHT:", 9.7, hy+0.1, 1.0, 0.25, size=7, bold=True, color=ACCENT2)
        add_text(s6, h.get("insight",""), 9.7, hy+0.38, 2.9, 0.42, size=9, color=TEXT, wrap=True)

    # ── SLIDE 7: Conversation ────────────────────────────────
    s7 = new_slide()
    add_eyebrow(s7, "BEYOND THE BRAND'S OWN CHANNELS")
    add_title(s7, "Where the Conversation Lives")
    add_divider(s7)
    convs = d.get("conversation",[])
    for i, ch in enumerate(convs):
        cy = 1.38 + i * 0.62
        add_rect(s7, 0.5, cy, 12.3, 0.54, CARD)
        add_text(s7, ch.get("platform",""), 1.1, cy+0.06, 2.8, 0.32, size=12, bold=True, color=WHITE)
        add_text(s7, ch.get("summary",""), 4.1, cy+0.06, 7.5, 0.42, size=10, color=TEXT, wrap=True)
        tc, tfc = tag_color(ch.get("tag",""))
        add_rect(s7, 11.15, cy+0.14, 1.25, 0.24, tc)
        add_text(s7, ch.get("tag",""), 11.15, cy+0.14, 1.25, 0.24, size=7, bold=True, color=tfc, align=PP_ALIGN.CENTER)

    # ── SLIDE 8: Ratings ─────────────────────────────────────
    s8 = new_slide()
    add_eyebrow(s8, "THE VERIFIED NUMBERS")
    add_title(s8, "Ratings Scorecard")
    add_divider(s8)
    ratings = d.get("ratings",[])
    for i, r in enumerate(ratings):
        ry = 1.38 + i * 0.68
        add_rect(s8, 0.5, ry, 12.3, 0.58, CARD)
        add_text(s8, r.get("platform",""), 1.1, ry+0.08, 3.5, 0.38, size=12, bold=True, color=WHITE)
        raw_r = str(r.get("rating","4.0")); import re as _re; _nums = _re.findall(r"[0-9]+.?[0-9]*", raw_r); rv = float(_nums[0]) if _nums else 4.0
        rc = GREEN if rv >= 4.5 else (AMBER if rv >= 4.0 else RED)
        add_text(s8, r.get("rating",""), 5.1, ry+0.05, 1.2, 0.5, size=22, bold=True, color=rc, align=PP_ALIGN.CENTER)
        add_text(s8, r.get("reviews",""), 6.5, ry+0.15, 4, 0.3, size=10, color=MUTED)
        tc, tfc = tag_color(r.get("tag",""))
        add_rect(s8, 11.0, ry+0.16, 1.4, 0.24, tc)
        add_text(s8, r.get("tag",""), 11.0, ry+0.16, 1.4, 0.24, size=7, bold=True, color=tfc, align=PP_ALIGN.CENTER)

    # ── SLIDE 9: Competitors ─────────────────────────────────
    s9 = new_slide()
    add_eyebrow(s9, "THE COMPETITIVE MAP")
    add_title(s9, "Competitor Landscape")
    add_divider(s9)
    comps = d.get("competitors",{}).get("list",[])
    for i, c in enumerate(comps):
        cy2 = 1.38 + i * 1.25
        add_rect(s9, 0.5, cy2, 12.3, 1.1, CARD)
        add_rect(s9, 0.6, cy2+0.22, 0.55, 0.55, RED)
        add_text(s9, str(i+1), 0.6, cy2+0.22, 0.55, 0.55, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s9, c.get("name",""), 1.28, cy2+0.1, 3, 0.38, size=15, bold=True, color=WHITE)
        add_text(s9, c.get("position",""), 1.28, cy2+0.5, 3, 0.3, size=10, italic=True, color=MUTED)
        add_text(s9, "STRENGTH:", 4.6, cy2+0.1, 1.1, 0.25, size=7, bold=True, color=MUTED)
        add_text(s9, c.get("strength",""), 5.8, cy2+0.08, 2.9, 0.38, size=10, color=TEXT, wrap=True)
        add_text(s9, "GAP TO OWN:", 8.8, cy2+0.1, 1.1, 0.25, size=7, bold=True, color=ACCENT2)
        add_text(s9, c.get("gap",""), 8.8, cy2+0.38, 3.8, 0.55, size=10, color=TEXT, wrap=True)
    key_gap = d.get("competitors",{}).get("key_gap","")
    if key_gap:
        add_rect(s9, 0.5, 5.15, 12.3, 0.65, rgb(0x0A,0x0F,0x2A))
        add_rect(s9, 0.5, 5.15, 0.07, 0.65, ACCENT)
        add_text(s9, f"KEY UNCONTESTED LANE — {key_gap}", 0.7, 5.22, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE 10: Voice of Customer ──────────────────────────
    s10 = new_slide()
    add_eyebrow(s10, "IN THEIR OWN WORDS")
    add_title(s10, "Voice of Customer")
    add_divider(s10)
    voc = d.get("voice_of_customer",{})
    add_rect(s10, 0.5, 1.38, 6.0, 4.0, CARD)
    add_rect(s10, 6.8, 1.38, 6.0, 4.0, CARD)
    add_rect(s10, 0.5, 1.38, 6.0, 0.45, GREEN)
    add_rect(s10, 6.8, 1.38, 6.0, 0.45, RED)
    add_text(s10, "✓  WHAT THEY LOVE", 0.5, 1.38, 6.0, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s10, "⚠  WHAT FRUSTRATES", 6.8, 1.38, 6.0, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, love in enumerate(voc.get("loves",[])):
        ly = 1.95 + i * 0.85
        add_text(s10, f"• {love.get('theme','')}", 0.65, ly, 5.7, 0.28, size=11, bold=True, color=WHITE)
        add_text(s10, love.get("quote",""), 0.65, ly+0.28, 5.7, 0.42, size=10, italic=True, color=TEXT, wrap=True)
    for i, fr in enumerate(voc.get("frustrations",[])):
        fy = 1.95 + i * 0.85
        add_text(s10, f"• {fr.get('theme','')}", 6.95, fy, 5.7, 0.28, size=11, bold=True, color=WHITE)
        add_text(s10, fr.get("quote",""), 6.95, fy+0.28, 5.7, 0.42, size=10, italic=True, color=TEXT, wrap=True)

    # ── SLIDE 11: Financial ──────────────────────────────────
    s11 = new_slide()
    add_eyebrow(s11, "COMPANY FINANCIALS & FUNDING")
    add_title(s11, "Financial Profile")
    add_divider(s11)
    fin = d.get("financial",{})
    fin_items = [("Total Funding",fin.get("total_funding","N/A")),("Last Round",fin.get("last_round","N/A")),
                 ("Revenue",fin.get("revenue","N/A")),("Valuation",fin.get("valuation","N/A")),
                 ("Employees",fin.get("employees","N/A")),("Founded",fin.get("founded","N/A"))]
    for i, (lbl, val) in enumerate(fin_items):
        col2, row = i % 3, i // 3
        fx = 0.5 + col2 * 4.15
        fy = 1.38 + row * 1.1
        add_rect(s11, fx, fy, 3.9, 0.95, CARD)
        add_text(s11, lbl, fx+0.15, fy+0.1, 3.6, 0.28, size=8, bold=True, color=MUTED)
        add_text(s11, val, fx+0.15, fy+0.38, 3.6, 0.45, size=13, bold=True, color=WHITE, wrap=True)
    investors = fin.get("investors",[])
    if any(investors):
        add_rect(s11, 0.5, 3.72, 12.3, 0.65, rgb(0x0A,0x1A,0x10))
        add_text(s11, "KEY INVESTORS", 0.65, 3.79, 2, 0.25, size=8, bold=True, color=MUTED)
        add_text(s11, "  ·  ".join(str(x) for x in investors if x), 2.8, 3.79, 9.7, 0.42, size=12, color=WHITE)
    if fin.get("hq"):
        add_rect(s11, 0.5, 4.55, 12.3, 0.55, rgb(0x0A,0x1A,0x10))
        add_rect(s11, 0.5, 4.55, 0.07, 0.55, GREEN)
        add_text(s11, f"HQ: {fin.get('hq','')}", 0.7, 4.62, 11.9, 0.4, size=11, color=WHITE)

    # ── SLIDE 12: Legal ──────────────────────────────────────
    s12 = new_slide()
    add_eyebrow(s12, "COMPANY REGISTRATION & LEGAL")
    add_title(s12, "Legal & Registration Profile")
    add_divider(s12)
    leg = d.get("company_legal",{})
    leg_items = [("CIN Number",leg.get("cin","N/A")),("Registered Name",leg.get("registered_name","N/A")),
                 ("Incorporation Date",leg.get("incorporation_date","N/A")),("MCA Status",leg.get("mca_status","N/A")),
                 ("DPIIT Recognition",leg.get("dpiit_status","N/A")),("Startup India",leg.get("startup_india","N/A"))]
    for i, (lbl, val) in enumerate(leg_items):
        col2, row = i % 2, i // 2
        lx = 0.5 + col2 * 6.3
        ly = 1.38 + row * 1.1
        add_rect(s12, lx, ly, 6.0, 0.95, CARD)
        add_text(s12, lbl, lx+0.15, ly+0.1, 5.7, 0.28, size=8, bold=True, color=MUTED)
        add_text(s12, val, lx+0.15, ly+0.4, 5.7, 0.42, size=13, bold=True, color=WHITE, wrap=True)
    dirs = leg.get("directors",[])
    if any(dirs):
        add_rect(s12, 0.5, 4.72, 12.3, 0.55, CARD)
        add_text(s12, f"DIRECTORS:  {'  ·  '.join(str(x) for x in dirs if x)}", 0.65, 4.79, 11.9, 0.4, size=11, color=WHITE)

    # ── SLIDE 13: YouTube ────────────────────────────────────
    s13 = new_slide()
    add_eyebrow(s13, "YOUTUBE CHANNEL INTELLIGENCE")
    add_title(s13, "YouTube Analytics")
    add_divider(s13)
    yt = d.get("youtube_real",{})
    yt_metrics = [
        ("Subscribers", yt_real["subscribers"] if yt_real else yt.get("subscribers","N/A")),
        ("Total Views",  yt_real["views"]       if yt_real else yt.get("total_views","N/A")),
        ("Videos",       yt_real["videos"]      if yt_real else yt.get("video_count","N/A")),
    ]
    for i, (lbl, val) in enumerate(yt_metrics):
        yx = 0.5 + i * 4.15
        add_rect(s13, yx, 1.38, 3.9, 1.7, CARD)
        add_text(s13, val, yx+0.15, 2.0, 3.6, 0.65, size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(s13, lbl, yx+0.15, 2.68, 3.6, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    if yt.get("top_content"):
        add_rect(s13, 0.5, 3.22, 12.3, 0.72, CARD)
        add_text(s13, "TOP CONTENT TYPE", 0.65, 3.3, 2.5, 0.25, size=8, bold=True, color=MUTED)
        add_text(s13, yt.get("top_content",""), 3.3, 3.29, 9.3, 0.5, size=11, color=TEXT, wrap=True)
    if yt.get("gap"):
        add_rect(s13, 0.5, 4.1, 12.3, 0.65, rgb(0x2A,0x0A,0x0A))
        add_rect(s13, 0.5, 4.1, 0.07, 0.65, RED)
        add_text(s13, f"GAP — {yt.get('gap','')}", 0.7, 4.17, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE 14: Gaps ───────────────────────────────────────
    s14 = new_slide()
    add_eyebrow(s14, "BLIND SPOTS TO CLOSE")
    add_title(s14, "Gaps, Risks & Whitespace")
    add_divider(s14)
    gaps = d.get("gaps",[])
    half = (len(gaps)+1)//2
    for i, g in enumerate(gaps):
        col2 = 0 if i < half else 1
        row = i if i < half else i - half
        gx = 0.5 + col2 * 6.3
        gy = 1.38 + row * 1.22
        add_rect(s14, gx, gy, 6.0, 1.08, CARD)
        add_text(s14, g.get("title",""), gx+0.15, gy+0.1, 4, 0.38, size=13, bold=True, color=WHITE)
        tc, tfc = tag_color(g.get("priority","MEDIUM"))
        add_rect(s14, gx+4.4, gy+0.1, 1.4, 0.26, tc)
        add_text(s14, g.get("priority",""), gx+4.4, gy+0.1, 1.4, 0.26, size=7, bold=True, color=tfc, align=PP_ALIGN.CENTER)
        add_text(s14, g.get("description",""), gx+0.15, gy+0.52, 5.7, 0.48, size=10, color=TEXT, wrap=True)

    # ── SLIDE 15: Recommendations ────────────────────────────
    s15 = new_slide()
    add_eyebrow(s15, "WHERE TO ACT FIRST")
    add_title(s15, "Priority Recommendations")
    add_divider(s15)
    recs = d.get("recommendations",[])
    for i, r in enumerate(recs):
        ry2 = 1.38 + i * 0.75
        add_rect(s15, 0.5, ry2, 12.3, 0.67, CARD)
        add_text(s15, r.get("number",str(i+1).zfill(2)), 0.65, ry2+0.1, 0.55, 0.48, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s15, r.get("title",""), 1.32, ry2+0.06, 3, 0.35, size=13, bold=True, color=WHITE)
        add_text(s15, r.get("description",""), 4.5, ry2+0.1, 8.0, 0.48, size=10, color=TEXT, wrap=True)

    # ── SLIDE: SEMrush Deep Dive ─────────────────────────────
    sSEM = new_slide()
    add_eyebrow(sSEM, "WEBSITE TRAFFIC & SEO INTELLIGENCE — SEMRUSH")
    add_title(sSEM, "SEMrush Analysis")
    add_divider(sSEM)
    sem = d.get("semrush_data", {})
    sem_metrics = [
        ("Monthly Visits",    sem.get("monthly_visits","N/A (est.)"),   "Organic + Paid combined"),
        ("Domain Rating",     sem.get("domain_rating","N/A (est.)"),    "Authority score out of 100"),
        ("Organic Keywords",  sem.get("organic_keywords","N/A (est.)"), "Keywords ranking on Google"),
        ("Total Backlinks",   sem.get("backlinks","N/A (est.)"),        "Referring domains count"),
    ]
    for i,(lbl,val,note) in enumerate(sem_metrics):
        col2, row = i%2, i//2
        sx = 0.5 + col2*6.3
        sy = 1.38 + row*1.1
        add_rect(sSEM, sx, sy, 6.0, 0.95, CARD)
        add_text(sSEM, lbl,  sx+0.15, sy+0.08, 5.7, 0.28, size=8, bold=True, color=MUTED)
        add_text(sSEM, val,  sx+0.15, sy+0.36, 5.7, 0.42, size=16, bold=True, color=ACCENT2)
        add_text(sSEM, note, sx+0.15, sy+0.72, 5.7, 0.22, size=9, color=MUTED)
    # Top keywords
    kws = sem.get("top_keywords", [])
    if kws:
        add_text(sSEM, "TOP ORGANIC KEYWORDS", 0.5, 3.65, 4, 0.28, size=8, bold=True, color=MUTED)
        for i, kw in enumerate(kws[:5]):
            kx = 0.5 + i*2.5
            add_rect(sSEM, kx, 4.0, 2.3, 0.38, rgb(0x1A,0x1A,0x2E))
            add_text(sSEM, str(kw), kx, 4.0, 2.3, 0.38, size=10, color=ACCENT2, align=PP_ALIGN.CENTER)
    # Traffic trend
    trend = sem.get("traffic_trend","")
    if trend:
        add_rect(sSEM, 0.5, 4.6, 12.3, 0.65, rgb(0x0A,0x1A,0x10))
        add_rect(sSEM, 0.5, 4.6, 0.07, 0.65, ACCENT2)
        add_text(sSEM, f"TRAFFIC TREND — {trend}", 0.7, 4.67, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE: Meta Ad Library Deep Dive ─────────────────────
    sMETA = new_slide()
    add_eyebrow(sMETA, "META AD LIBRARY — FACEBOOK & INSTAGRAM ADS")
    add_title(sMETA, "Meta Ad Library Analysis")
    add_divider(sMETA)
    meta = d.get("meta_ads_data", {})
    # Stats row
    meta_stats = [
        ("Active Ads Found",  meta.get("total_ads","N/A")),
        ("Ad Formats",        meta.get("primary_format","Video + Image")),
        ("Primary CTA",       meta.get("primary_cta","N/A")),
        ("Avg Run Duration",  meta.get("avg_duration","N/A")),
    ]
    for i,(lbl,val) in enumerate(meta_stats):
        mx = 0.5 + i*3.1
        add_rect(sMETA, mx, 1.38, 2.95, 0.95, CARD)
        add_text(sMETA, lbl, mx+0.12, 1.45, 2.7, 0.28, size=8, bold=True, color=MUTED)
        add_text(sMETA, val, mx+0.12, 1.72, 2.7, 0.52, size=12, bold=True, color=WHITE, wrap=True)
    # Ad format bars
    add_text(sMETA, "AD FORMAT BREAKDOWN", 0.5, 2.55, 6, 0.28, size=8, bold=True, color=MUTED)
    formats = meta.get("formats", [{"type":"Video Ads","pct":70},{"type":"Static Image","pct":20},{"type":"Carousel","pct":10}])
    for i,f in enumerate(formats):
        fy = 2.9 + i*0.7
        add_text(sMETA, f.get("type",""), 0.5, fy, 2.5, 0.28, size=11, bold=True, color=WHITE)
        pct = f.get("pct",0)
        add_rect(sMETA, 3.1, fy+0.04, 5.5, 0.3, rgb(0x2A,0x2A,0x3A))
        add_rect(sMETA, 3.1, fy+0.04, 5.5*(pct/100), 0.3, ACCENT)
        add_text(sMETA, f"{pct}%", 8.8, fy+0.04, 0.8, 0.3, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    # Top hooks
    hooks_raw = meta.get("top_hooks", [])
    if hooks_raw:
        add_text(sMETA, "TOP AD HOOKS RUNNING", 0.5, 5.0, 6, 0.28, size=8, bold=True, color=MUTED)
        for i,h in enumerate(hooks_raw[:2]):
            hy = 5.35 + i*0.55
            add_rect(sMETA, 0.5, hy, 12.3, 0.45, rgb(0x1A,0x1A,0x2E))
            add_text(sMETA, f"→  {h}", 0.7, hy+0.08, 11.9, 0.3, size=10, color=TEXT)
    # Gap
    gap = meta.get("gap","")
    if gap:
        add_rect(sMETA, 0.5, 6.1, 12.3, 0.65, rgb(0x2A,0x0A,0x0A))
        add_rect(sMETA, 0.5, 6.1, 0.07, 0.65, RED)
        add_text(sMETA, f"CRITICAL GAP — {gap}", 0.7, 6.17, 11.9, 0.52, size=10, color=TEXT)

    # ── SLIDE: Google Ad Transparency Deep Dive ──────────────
    sGADS = new_slide()
    add_eyebrow(sGADS, "GOOGLE AD TRANSPARENCY — SEARCH, DISPLAY & YOUTUBE ADS")
    add_title(sGADS, "Google Ads Intelligence")
    add_divider(sGADS)
    gads = d.get("google_ads_data", {})
    # Stats
    gads_stats = [
        ("Active Google Ads",    gads.get("total_ads","N/A")),
        ("Ad Types Running",     gads.get("ad_types","Search + Display")),
        ("Primary Keywords",     gads.get("primary_keywords","N/A")),
        ("Est. Monthly Spend",   gads.get("est_spend","N/A")),
    ]
    for i,(lbl,val) in enumerate(gads_stats):
        gx2 = 0.5 + i*3.1
        add_rect(sGADS, gx2, 1.38, 2.95, 0.95, CARD)
        add_text(sGADS, lbl, gx2+0.12, 1.45, 2.7, 0.28, size=8, bold=True, color=MUTED)
        add_text(sGADS, val, gx2+0.12, 1.72, 2.7, 0.52, size=12, bold=True, color=WHITE, wrap=True)
    # Ad types breakdown
    add_text(sGADS, "AD TYPES BREAKDOWN", 0.5, 2.55, 6, 0.28, size=8, bold=True, color=MUTED)
    ad_types = gads.get("ad_type_breakdown", [
        {"type":"Search Ads","pct":50,"note":"Text ads on Google Search"},
        {"type":"Display Ads","pct":30,"note":"Banner ads across websites"},
        {"type":"YouTube Ads","pct":20,"note":"Video ads on YouTube"},
    ])
    for i,at in enumerate(ad_types):
        ty = 2.9 + i*0.82
        add_rect(sGADS, 0.5, ty, 12.3, 0.7, CARD)
        add_text(sGADS, at.get("type",""), 0.65, ty+0.08, 2.5, 0.32, size=12, bold=True, color=WHITE)
        pct2 = at.get("pct",0)
        add_rect(sGADS, 3.2, ty+0.18, 5.5, 0.28, rgb(0x2A,0x2A,0x3A))
        add_rect(sGADS, 3.2, ty+0.18, 5.5*(pct2/100), 0.28, ACCENT2)
        add_text(sGADS, f"{pct2}%", 8.85, ty+0.14, 0.8, 0.32, size=12, bold=True, color=ACCENT2, align=PP_ALIGN.RIGHT)
        add_text(sGADS, at.get("note",""), 9.8, ty+0.14, 2.8, 0.32, size=9, color=MUTED, wrap=True)
    # Key insight
    insight = gads.get("insight","")
    if insight:
        add_rect(sGADS, 0.5, 5.45, 12.3, 0.65, rgb(0x0A,0x0F,0x2A))
        add_rect(sGADS, 0.5, 5.45, 0.07, 0.65, ACCENT)
        add_text(sGADS, f"KEY INSIGHT — {insight}", 0.7, 5.52, 11.9, 0.52, size=10, color=TEXT)

    # ── END SLIDE ────────────────────────────────────────────
    end = prs.slides.add_slide(blank)
    slide_bg(end)
    add_rect(end, 0, 0, 0.12, 7.5, ACCENT2)
    add_text(end, "BRAND INTELLIGENCE", 0.5, 1.6, 12, 0.4, size=10, bold=True, color=MUTED)
    add_text(end, d.get("brandName","Brand"), 0.5, 2.1, 12, 1.2, size=52, bold=True, color=WHITE)
    add_rect(end, 0.5, 3.5, 2.5, 0.04, ACCENT2)
    add_text(end, "Automated public-web scan  ·  June 2026  ·  PREPARED FOR INTERNAL REVIEW",
             0.5, 6.8, 12, 0.35, size=9, color=MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ── API Route ─────────────────────────────────────────────────
@app.route("/generate", methods=["POST"])
def generate():
    data = request.json
    brand = data.get("brand","").strip()
    website = data.get("website","").strip()
    if not brand:
        return jsonify({"error":"Brand name required"}), 400

    try:
        print(f"[1/4] Researching {brand} | website: {website}")
        raw = deep_research(brand, website)

        print(f"[2/4] Getting YouTube data...")
        yt_real = get_youtube(brand)

        print(f"[3/4] Claude analysis...")
        report = claude_analysis(brand, raw, yt_real)

        print(f"[4/4] Building PPT...")
        ppt_buf = build_ppt(report, yt_real)

        safe = brand.replace(" ","-").replace("/","")
        return send_file(
            ppt_buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"{safe}-Intelligence-Report.pptx"
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/health")
def health():
    return jsonify({"status":"ok","message":"Brand PPT Generator running"})

@app.route("/")
def home():
    return send_file("index.html")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
